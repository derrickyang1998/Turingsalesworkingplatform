#!/usr/bin/env python3
"""Generate editable light-theme PPTX decks from TuringMarket proposal outlines."""
import json
import os
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PAGE_W = 13.333
PAGE_H = 7.5

INK = RGBColor(17, 24, 39)
BODY = RGBColor(55, 65, 81)
MUTED = RGBColor(107, 114, 128)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(248, 250, 252)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(229, 231, 235)
SOFT_BLUE = RGBColor(238, 242, 255)
SOFT_PURPLE = RGBColor(245, 243, 255)
SOFT_MINT = RGBColor(236, 253, 245)
ACCENT = RGBColor(79, 70, 229)
ACCENT_2 = RGBColor(124, 58, 237)
GREEN = RGBColor(16, 185, 129)


def clean(value):
    return str(value or "").replace("\r", "").strip()


def split_metric(text):
    text = clean(text)
    for sep in ("：", ":"):
        if sep in text:
            left, right = text.split(sep, 1)
            return clean(left), clean(right)
    return text[:24] or "策略要点", text if len(text) > 24 else ""


def metric_lead(text):
    match = re.search(r"(\d+%|\$[0-9][0-9,Kk+\-– ]*|[0-9]+[KkWw万+]*)", clean(text))
    return match.group(1) if match else "●"


def compact(text, limit):
    text = re.sub(r"\s+", " ", clean(text))
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 1)].rstrip() + "…"


def infer_brand(title):
    title = clean(title)
    return title.split()[0] if title else "CLIENT"


def normalize_points(points):
    if isinstance(points, str):
        points = [p.strip() for p in re.split(r"[;\n；]+", points) if p.strip()]
    if not isinstance(points, list):
        return []
    return [clean(p) for p in points if clean(p)]


def normalize(data):
    if not isinstance(data, dict):
        data = {}
    outer_research = data.get("research")
    outer_demand = data.get("demand") if isinstance(data.get("demand"), dict) else {}
    if "outline" in data:
        data = data.get("outline") or {}
        if outer_research and not data.get("research"):
            data["research"] = outer_research

    sections = data.get("sections") or []
    normalized = []
    for sec in sections:
        if not isinstance(sec, dict):
            continue
        normalized.append({
            "title": clean(sec.get("title") or "Proposal Slide"),
            "type": clean(sec.get("type") or "content"),
            "points": normalize_points(sec.get("points") or sec.get("items") or []),
            "note": clean(sec.get("note")),
        })

    title = clean(data.get("title") or "TuringMarket Influencer Marketing Proposal")
    return {
        "title": title,
        "subtitle": clean(data.get("subtitle") or "Client proposal deck"),
        "brand": clean(data.get("brand") or outer_demand.get("brand") or outer_demand.get("company") or infer_brand(title)),
        "product": clean(data.get("product") or outer_demand.get("product")),
        "research": data.get("research") or {},
        "materials": data.get("materials") or [],
        "sections": normalized,
    }


def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1)
    return shape


def add_text(slide, text, x, y, w, h, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = clean(text)
    p.alignment = align
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_pill(slide, text, x, y, w, fill=SOFT_BLUE, color=ACCENT):
    add_rect(slide, x, y, w, 0.34, fill, fill)
    add_text(slide, text, x + 0.08, y + 0.07, w - 0.16, 0.18, 8.5, True, color, PP_ALIGN.CENTER)


def add_decor(slide):
    add_rect(slide, -0.35, -0.25, 3.4, 1.05, SOFT_BLUE, SOFT_BLUE, True)
    add_rect(slide, 10.95, 0.18, 2.4, 0.72, SOFT_MINT, SOFT_MINT, True)
    add_rect(slide, 11.7, 6.8, 1.4, 0.24, SOFT_PURPLE, SOFT_PURPLE, True)


def add_header(slide, title, subtitle="", page_num=1):
    add_decor(slide)
    add_text(slide, "TURINGMARKET", 0.62, 0.36, 2.1, 0.2, 8.5, True, ACCENT)
    add_text(slide, str(page_num).zfill(2), 12.28, 0.36, 0.45, 0.2, 8.5, True, MUTED, PP_ALIGN.RIGHT)
    add_text(slide, compact(title, 52), 0.62, 0.75, 11.35, 0.62, 28, True, INK)
    if subtitle:
        add_text(slide, compact(subtitle, 116), 0.65, 1.34, 10.8, 0.26, 10.5, False, MUTED)
    add_rect(slide, 0.62, 1.78, 0.68, 0.05, ACCENT, ACCENT, False)
    add_rect(slide, 1.34, 1.78, 10.8, 0.02, LINE, LINE, False)


def add_footer(slide, page_num, total):
    add_text(slide, "TuringMarket 图灵集市 | 海外红人营销提案", 0.62, 7.05, 5.8, 0.2, 8.5, False, MUTED)
    add_text(slide, f"{str(page_num).zfill(2)} / {str(total).zfill(2)}", 11.6, 7.05, 1.0, 0.2, 8.5, True, MUTED, PP_ALIGN.RIGHT)
    add_rect(slide, 0.62, 6.92, 12.1, 0.02, LINE, LINE, False)
    if total:
        add_rect(slide, 0.62, 6.92, 12.1 * page_num / total, 0.035, ACCENT, ACCENT, False)


def section_label(sec):
    text = clean(sec.get("title"))
    slide_type = sec.get("type", "content")
    if slide_type == "research" or "调研" in text:
        return "MARKET RESEARCH"
    if slide_type == "sources" or "来源" in text:
        return "SOURCES"
    if slide_type == "timeline" or "排期" in text:
        return "TIMELINE"
    if slide_type in ("stats", "kpi") or re.search(r"预算|平台|KPI|指标|budget|platform", text, re.I):
        return "BUDGET & KPI"
    if slide_type in ("next", "closing"):
        return "NEXT STEPS"
    if re.search(r"内容|脚本|创意", text):
        return "CONTENT STRATEGY"
    if re.search(r"达人|红人", text):
        return "INFLUENCER MATRIX"
    return "STRATEGY"


def cover_slide(prs, data, sec, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_rect(slide, 0.0, 0.0, PAGE_W, PAGE_H, WHITE, WHITE, False)
    add_rect(slide, 0.52, 0.52, 12.28, 6.45, BG, LINE, True)
    add_rect(slide, 0.92, 0.94, 3.0, 0.38, SOFT_BLUE, SOFT_BLUE, True)
    add_text(slide, "INFLUENCER MARKETING STRATEGY", 1.05, 1.05, 2.74, 0.14, 7.5, True, ACCENT, PP_ALIGN.CENTER)
    title = sec.get("title") or data["title"]
    subtitle = sec.get("points", [data["subtitle"]])[0] if sec.get("points") else data["subtitle"]
    add_text(slide, compact(title, 86), 0.95, 1.72, 7.75, 1.35, 37, True, INK)
    add_text(slide, compact(subtitle, 150), 0.98, 3.25, 7.2, 0.55, 15, False, MUTED)
    chips = ["策略规划", "达人矩阵", "内容脚本", "数据复盘"]
    for i, chip in enumerate(chips):
        add_pill(slide, chip, 0.98 + i * 1.25, 4.24, 1.08)
    add_rect(slide, 9.08, 1.38, 2.85, 4.35, WHITE, LINE, True)
    add_text(slide, "CLIENT DECK", 9.42, 1.78, 2.15, 0.24, 10, True, ACCENT, PP_ALIGN.CENTER)
    for idx, label in enumerate(["需求理解", "市场判断", "执行落地", "复盘增长"]):
        y = 2.35 + idx * 0.66
        add_rect(slide, 9.55, y, 0.28, 0.28, ACCENT if idx == 0 else SOFT_BLUE, ACCENT if idx == 0 else SOFT_BLUE, True)
        add_text(slide, label, 10.0, y - 0.01, 1.35, 0.25, 11.5, True, BODY)
    add_text(slide, "Presented by TuringMarket 图灵集市", 0.98, 6.3, 4.4, 0.2, 9.5, False, MUTED)
    add_footer(slide, page_num, total)


def content_slide(prs, data, sec, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_header(slide, sec["title"], sec.get("note") or data["subtitle"], page_num)
    points = sec.get("points", [])[:6] or ["待补充具体内容"]
    positions = [(0.72, 2.08), (4.72, 2.08), (8.72, 2.08), (0.72, 4.35), (4.72, 4.35), (8.72, 4.35)]
    for idx, (point, (x, y)) in enumerate(zip(points, positions)):
        label, body = split_metric(point)
        add_rect(slide, x, y, 3.52, 1.78, PANEL, LINE, True)
        add_text(slide, f"{idx + 1:02d}", x + 0.24, y + 0.22, 0.38, 0.18, 8.5, True, ACCENT)
        add_text(slide, compact(label, 36), x + 0.24, y + 0.49, 2.95, 0.34, 13.2, True, INK)
        add_text(slide, compact(body or point, 132), x + 0.24, y + 0.92, 2.95, 0.56, 10.3, False, MUTED)
    add_footer(slide, page_num, total)


def stats_slide(prs, data, sec, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_header(slide, sec["title"], sec.get("note") or data["subtitle"], page_num)
    points = sec.get("points", [])[:6] or ["核心指标: 待补充"]
    positions = [(0.72, 2.08), (3.88, 2.08), (7.04, 2.08), (10.2, 2.08)]
    for idx, point in enumerate(points[:4]):
        label, body = split_metric(point)
        x, y = positions[idx]
        add_rect(slide, x, y, 2.72, 1.7, PANEL, LINE, True)
        add_rect(slide, x, y, 2.72, 0.08, ACCENT if idx % 2 == 0 else ACCENT_2, ACCENT if idx % 2 == 0 else ACCENT_2, False)
        add_text(slide, metric_lead(label + " " + body), x + 0.22, y + 0.34, 1.2, 0.36, 22, True, ACCENT)
        add_text(slide, compact(label, 28), x + 0.22, y + 0.82, 2.2, 0.26, 12, True, INK)
        add_text(slide, compact(body or point, 72), x + 0.22, y + 1.16, 2.24, 0.28, 9.2, False, MUTED)
    add_rect(slide, 0.72, 4.25, 11.9, 1.72, PANEL, LINE, True)
    add_text(slide, "执行拆解", 1.02, 4.52, 1.2, 0.2, 10, True, ACCENT)
    for idx, point in enumerate(points[:5]):
        label, body = split_metric(point)
        y = 4.85 + idx * 0.23
        add_text(slide, compact(label, 20), 1.02, y, 2.0, 0.18, 8.7, True, BODY)
        add_text(slide, compact(body or point, 96), 3.05, y, 8.75, 0.18, 8.5, False, MUTED)
    add_footer(slide, page_num, total)


def timeline_slide(prs, data, sec, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_header(slide, sec["title"], sec.get("note") or data["subtitle"], page_num)
    points = sec.get("points", [])[:5] or ["阶段|时间|核心动作|交付物"]
    y = 2.1
    for idx, point in enumerate(points):
        parts = [clean(p) for p in point.split("|")]
        phase = parts[0] if len(parts) > 0 and parts[0] else "阶段 " + str(idx + 1)
        time = parts[1] if len(parts) > 1 and parts[1] else "待确认"
        action = parts[2] if len(parts) > 2 and parts[2] else point
        output = parts[3] if len(parts) > 3 and parts[3] else "阶段交付物"
        add_rect(slide, 0.72, y, 11.9, 0.68, PANEL, LINE, True)
        add_rect(slide, 0.96, y + 0.17, 0.34, 0.34, ACCENT, ACCENT, True)
        add_text(slide, str(idx + 1), 1.06, y + 0.25, 0.13, 0.1, 8.5, True, WHITE, PP_ALIGN.CENTER)
        add_text(slide, compact(phase, 24), 1.5, y + 0.16, 1.65, 0.22, 10.5, True, ACCENT)
        add_text(slide, compact(time, 18), 3.15, y + 0.16, 1.25, 0.22, 9.5, True, MUTED)
        add_text(slide, compact(action, 72), 4.55, y + 0.15, 3.65, 0.24, 9.5, False, BODY)
        add_text(slide, compact(output, 54), 8.45, y + 0.15, 3.55, 0.24, 9.5, False, MUTED)
        y += 0.84
    add_footer(slide, page_num, total)


def research_slide(prs, data, sec, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_header(slide, sec["title"], sec.get("note") or data["subtitle"], page_num)
    points = sec.get("points", [])[:6] or ["调研状态: 未读取到可用在线来源"]
    positions = [(0.72, 2.08), (6.75, 2.08), (0.72, 3.56), (6.75, 3.56), (0.72, 5.04), (6.75, 5.04)]
    for idx, (point, (x, y)) in enumerate(zip(points, positions)):
        label, body = split_metric(point)
        add_rect(slide, x, y, 5.52, 1.08, PANEL, LINE, True)
        add_text(slide, "SIGNAL " + str(idx + 1).zfill(2), x + 0.22, y + 0.18, 1.15, 0.16, 8.2, True, ACCENT)
        add_text(slide, compact(label, 34), x + 1.42, y + 0.17, 1.95, 0.2, 10, True, INK)
        add_text(slide, compact(body or point, 102), x + 0.22, y + 0.52, 4.9, 0.34, 8.8, False, MUTED)
    add_footer(slide, page_num, total)


def sources_slide(prs, data, sec, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_header(slide, sec["title"], sec.get("note") or data["subtitle"], page_num)
    points = sec.get("points", [])[:8] or ["来源状态: 暂无可引用在线来源"]
    positions = [(0.72, 2.0), (6.75, 2.0), (0.72, 3.15), (6.75, 3.15), (0.72, 4.3), (6.75, 4.3), (0.72, 5.45), (6.75, 5.45)]
    for point, (x, y) in zip(points, positions):
        label, body = split_metric(point)
        add_rect(slide, x, y, 5.52, 0.86, PANEL, LINE, True)
        add_text(slide, compact(label, 42), x + 0.2, y + 0.14, 5.05, 0.18, 8.8, True, ACCENT)
        add_text(slide, compact(body or point, 112), x + 0.2, y + 0.43, 5.0, 0.22, 7.8, False, MUTED)
    add_footer(slide, page_num, total)


def materials_slide(prs, data, sec, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_header(slide, sec["title"], sec.get("note") or data["subtitle"], page_num)
    points = sec.get("points", [])[:6] or ["未上传补充材料"]
    positions = [(0.72, 2.08), (6.75, 2.08), (0.72, 3.56), (6.75, 3.56), (0.72, 5.04), (6.75, 5.04)]
    for point, (x, y) in zip(points, positions):
        label, body = split_metric(point)
        add_rect(slide, x, y, 5.52, 1.08, PANEL, LINE, True)
        add_pill(slide, "SOURCE MATERIAL", x + 0.2, y + 0.17, 1.34, SOFT_BLUE, ACCENT)
        add_text(slide, compact(label, 38), x + 1.72, y + 0.19, 3.65, 0.18, 9.5, True, INK)
        add_text(slide, compact(body or point, 116), x + 0.2, y + 0.55, 5.0, 0.28, 8.4, False, MUTED)
    add_footer(slide, page_num, total)


def closing_slide(prs, data, sec, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_rect(slide, 0.52, 0.52, 12.28, 6.45, BG, LINE, True)
    title = sec.get("title") or "让品牌在目标市场建立可复用增长资产"
    add_text(slide, compact(title, 64), 1.15, 1.55, 7.9, 1.05, 34, True, INK)
    add_text(slide, compact(data.get("subtitle") or "策略 · 达人 · 内容 · 数据闭环", 126), 1.18, 2.95, 7.0, 0.42, 14, False, MUTED)
    points = sec.get("points", [])[:4] or ["确认预算与市场优先级", "确认达人筛选红线", "进入达人名单匹配与执行排期"]
    for idx, point in enumerate(points):
        add_pill(slide, compact(point, 20), 1.18 + (idx % 2) * 2.55, 3.78 + (idx // 2) * 0.52, 2.28)
    add_rect(slide, 9.35, 1.62, 2.32, 2.32, SOFT_BLUE, SOFT_BLUE, True)
    add_text(slide, "THANK\nYOU", 9.75, 2.2, 1.52, 0.78, 28, True, ACCENT, PP_ALIGN.CENTER)
    add_text(slide, "Presented by TuringMarket 图灵集市", 1.18, 6.18, 4.1, 0.2, 9.5, False, MUTED)
    add_footer(slide, page_num, total)


def next_slide(prs, data, sec, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_header(slide, sec["title"], sec.get("note") or data["subtitle"], page_num)
    points = sec.get("points", [])[:6] or ["确认下一步合作"]
    positions = [(0.72, 2.18), (4.72, 2.18), (8.72, 2.18), (0.72, 4.45), (4.72, 4.45), (8.72, 4.45)]
    for idx, (point, (x, y)) in enumerate(zip(points, positions)):
        label, body = split_metric(point)
        add_rect(slide, x, y, 3.52, 1.62, PANEL, LINE, True)
        add_rect(slide, x, y, 3.52, 0.48, SOFT_BLUE, SOFT_BLUE, True)
        add_text(slide, "NEXT " + str(idx + 1).zfill(2), x + 0.22, y + 0.15, 0.72, 0.12, 7.5, True, ACCENT)
        add_text(slide, compact(label, 34), x + 0.22, y + 0.64, 2.95, 0.32, 12.2, True, INK)
        add_text(slide, compact(body or point, 110), x + 0.22, y + 1.04, 2.95, 0.34, 9.2, False, MUTED)
    add_footer(slide, page_num, total)


def research_points_from_payload(research):
    bullets = research.get("bullets") or []
    points = [clean(item) for item in bullets if clean(item)]
    if not points:
        for idx, item in enumerate(research.get("sources") or []):
            title = clean(item.get("title") or item.get("url") or "")
            snippet = clean(item.get("snippet") or item.get("url") or "")
            if title:
                points.append(f"来源{idx + 1}: {title} - {snippet}")
    return points[:6]


def source_points_from_payload(research):
    points = []
    for idx, item in enumerate(research.get("sources") or []):
        title = clean(item.get("title") or item.get("url") or "")
        url = clean(item.get("url") or "")
        snippet = clean(item.get("snippet") or "")
        if title or url:
            points.append(f"来源{idx + 1}: {title} | {snippet or url}")
    return points[:8]


def apply_research_sections(data, sections):
    research = data.get("research") or {}
    if not isinstance(research, dict):
        return sections
    output = list(sections)
    has_research = any(sec.get("type") == "research" or "调研" in sec.get("title", "") for sec in output)
    has_sources = any(sec.get("type") == "sources" or "来源" in sec.get("title", "") for sec in output)
    research_points = research_points_from_payload(research)
    source_points = source_points_from_payload(research)
    if research_points and not has_research:
        output.insert(min(3, len(output)), {
            "title": "联网调研与市场信号",
            "type": "research",
            "points": research_points,
            "note": "生成前联网搜索结果已用于校准市场、竞品、平台和内容落地。",
        })
    if source_points and not has_sources:
        output.append({
            "title": "调研来源与引用口径",
            "type": "sources",
            "points": source_points,
            "note": "来源用于策略判断参考，最终执行前仍建议乙方二次核验。",
        })
    return output


def build_materials_section(data):
    materials = data.get("materials") or []
    if not materials:
        return None
    points = []
    for item in materials[:8]:
        if not isinstance(item, dict):
            continue
        name = clean(item.get("name") or "补充材料")
        state = "已提取正文" if item.get("parsed") else "仅文件信息"
        preview = clean(item.get("preview") or item.get("text") or "未提取到可读正文")
        points.append(f"{name}: {state} | {preview}")
    if not points:
        return None
    return {
        "title": "补充材料引用",
        "type": "materials",
        "points": points,
        "note": "以下内容来自上传的 PDF / Word / PPTX / Excel / 文本文件。",
    }


def build_section_list(data):
    sections = list(data.get("sections") or [])
    if not sections or sections[0].get("type") != "cover":
        sections.insert(0, {
            "title": data["title"],
            "type": "cover",
            "points": [data["subtitle"]],
            "note": "",
        })
    sections = apply_research_sections(data, sections)
    material_section = build_materials_section(data)
    if material_section:
        sections.append(material_section)
    sections.append({
        "title": "让 " + (data.get("brand") or infer_brand(data["title"])) + " 在目标市场建立可复用增长资产",
        "type": "closing",
        "points": ["策略 · 达人 · 内容 · 数据闭环", "可执行交付物", "HTMLPPT 与 PPTX 同版式输出"],
        "note": "",
    })
    return sections


def build_deck(data, out_path):
    prs = Presentation()
    prs.slide_width = Inches(PAGE_W)
    prs.slide_height = Inches(PAGE_H)
    sections = build_section_list(data)
    total = len(sections)

    for index, sec in enumerate(sections, start=1):
        slide_type = sec.get("type", "content")
        if slide_type == "cover" or index == 1:
            cover_slide(prs, data, sec, index, total)
        elif slide_type == "research":
            research_slide(prs, data, sec, index, total)
        elif slide_type == "sources":
            sources_slide(prs, data, sec, index, total)
        elif slide_type == "materials":
            materials_slide(prs, data, sec, index, total)
        elif slide_type in ("stats", "kpi") or re.search(r"预算|平台|KPI|指标|budget|platform", sec.get("title", ""), re.I):
            stats_slide(prs, data, sec, index, total)
        elif slide_type == "timeline" or re.search(r"排期|timeline|里程碑", sec.get("title", ""), re.I):
            timeline_slide(prs, data, sec, index, total)
        elif slide_type == "next":
            next_slide(prs, data, sec, index, total)
        elif slide_type == "closing":
            closing_slide(prs, data, sec, index, total)
        else:
            content_slide(prs, data, sec, index, total)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)


if __name__ == "__main__":
    data_path = sys.argv[1]
    out_path = sys.argv[2]
    with open(data_path, "r", encoding="utf-8-sig") as f:
        payload = json.load(f)
    build_deck(normalize(payload), out_path)
    print("OK")
